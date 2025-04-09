import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from io import BytesIO

# Fonction finale avec gestion valeurs vides
def remplacer_placeholder(slide, placeholder, nouvelle_valeur):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if placeholder in paragraph.text:
                    premier_run = paragraph.runs[0] if paragraph.runs else None
                    if premier_run:
                        style = {
                            'font_name': premier_run.font.name,
                            'font_size': premier_run.font.size,
                            'bold': premier_run.font.bold,
                            'italic': premier_run.font.italic,
                            'underline': premier_run.font.underline,
                            'color': premier_run.font.color.rgb if premier_run.font.color.type else None,
                        }
                    else:
                        style = None

                    remplacement = nouvelle_valeur if nouvelle_valeur and pd.notna(nouvelle_valeur) else ""

                    texte_modifie = paragraph.text.replace(placeholder, remplacement)

                    p = paragraph._element
                    for r in list(p):
                        p.remove(r)

                    nouveau_run = paragraph.add_run()
                    nouveau_run.text = texte_modifie

                    if style:
                        nouveau_run.font.name = style['font_name']
                        nouveau_run.font.size = style['font_size']
                        nouveau_run.font.bold = style['bold']
                        nouveau_run.font.italic = style['italic']
                        nouveau_run.font.underline = style['underline']
                        if style['color']:
                            nouveau_run.font.color.rgb = style['color']

st.title("🚀 Générateur PPTX Final (Placeholders invisibles si vide)")

fichier_excel = st.file_uploader("📥 Charge ton fichier Excel", type=["xlsx"])
fichier_pptx = st.file_uploader("📥 Charge ton modèle PowerPoint", type=["pptx"])

if fichier_excel and fichier_pptx:
    donnees = pd.read_excel(fichier_excel)
    prs = Presentation(fichier_pptx)

    placeholders = {
        '{{Titre_Residence}}': 'Nom de la residence',
        '{{Ville}}': 'Ville',
        '{{Details}}': 'details',
        '{{Maitre_Ouvrage}}': 'Maitre d\'ouvrage',
        '{{Assistant}}': 'Assistant',
        '{{Montant_Travaux}}': 'montant des travaux',
        '{{Type_Travaux}}': 'type de travaux',
        '{{Dates_Realisation}}': 'realisation',
        '{{Surface}}': 'Surface',
        '{{Travaux_Exterieurs}}': 'travaux exterieurs',
        '{{Travaux_Interieurs}}': 'travaux interieurs'
    }

    for slide, (_, projet) in zip(prs.slides, donnees.iterrows()):
        for placeholder, colonne in placeholders.items():
            valeur = projet[colonne]
            remplacer_placeholder(slide, placeholder, str(valeur) if pd.notna(valeur) else "")

    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)

    st.download_button(
        label="📤 Télécharger PPTX parfaitement stylisé",
        data=pptx_io,
        file_name="Presentation_Finalisee.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
